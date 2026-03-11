import os
import asyncio
import google.generativeai as genai
from pptx import Presentation
from docx import Document
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import Command
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.types import ReplyKeyboardMarkup, KeyboardButton, FSInputFile

# --- KONFIGURATSIYA ---
# Railway o'zgaruvchilarini o'qiymiz
TOKEN = os.environ.get("BOT_TOKEN")
GEMINI_KEY = os.environ.get("GEMINI_API_KEY")

# Gemini API ni sozlash
genai.configure(api_key=GEMINI_KEY)

# MUHIM: Modelni eng barqaror shaklda chaqiramiz
# Agar 'gemini-1.5-flash' ishlamasa, 'models/gemini-pro' ni sinab ko'radi
try:
    model = genai.GenerativeModel('models/gemini-1.5-flash')
except:
    model = genai.GenerativeModel('models/gemini-pro')

bot = Bot(token=TOKEN)
dp = Dispatcher()

class BotStates(StatesGroup):
    waiting_for_esse = State()
    waiting_for_pptx = State()

# --- ASOSIY MENYU ---
menu = ReplyKeyboardMarkup(keyboard=[
    [KeyboardButton(text="📝 Sifatli Esse (Word)")],
    [KeyboardButton(text="📊 Professional Taqdimot (PPTX)")]
], resize_keyboard=True)

@dp.message(Command("start"))
async def start(message: types.Message, state: FSMContext):
    await state.clear()
    await message.answer(
        "Salom! Bot muvaffaqiyatli yangilandi. 🚀\n"
        "Quyidagi xizmatlardan birini tanlang:", 
        reply_markup=menu
    )

# --- ESSE YARATISH ---
@dp.message(F.text == "📝 Sifatli Esse (Word)")
async def esse_req(message: types.Message, state: FSMContext):
    await state.set_state(BotStates.waiting_for_esse)
    await message.answer("📝 Esse mavzusini yuboring (Masalan: 'O'zbekiston madaniyati'):")

@dp.message(BotStates.waiting_for_esse)
async def handle_esse(message: types.Message, state: FSMContext):
    msg = await message.answer("🧠 AI akademik esse tayyorlamoqda, iltimos kuting...")
    try:
        # AI dan javob olish
        response = model.generate_content(f"{message.text} mavzusida o'zbek tilida batafsil akademik esse yoz.")
        
        # Word fayl yaratish
        doc = Document()
        doc.add_heading(message.text, 0)
        doc.add_paragraph(response.text)
        
        file_path = f"esse_{message.from_user.id}.docx"
        doc.save(file_path)
        
        # Faylni yuborish
        await message.answer_document(FSInputFile(file_path), caption=f"✅ '{message.text}' bo'yicha esse tayyor!")
        os.remove(file_path)
    except Exception as e:
        await message.answer(f"❌ Xatolik yuz berdi. API kalit yoki ulanishni qayta tekshiring.")
        print(f"ERROR: {e}")
    finally:
        await msg.delete()
        await state.clear()

# --- TAQDIMOT YARATISH ---
@dp.message(F.text == "📊 Professional Taqdimot (PPTX)")
async def pptx_req(message: types.Message, state: FSMContext):
    await state.set_state(BotStates.waiting_for_pptx)
    await message.answer("📊 Taqdimot mavzusini yuboring:")

@dp.message(BotStates.waiting_for_pptx)
async def handle_pptx(message: types.Message, state: FSMContext):
    msg = await message.answer("🎨 Taqdimot slaydlari tayyorlanmoqda...")
    try:
        prompt = f"'{message.text}' mavzusida 7 ta slayd uchun matn tayyorla. Har bir slaydni 'SLAYD:' so'zi bilan boshla."
        response = model.generate_content(prompt)
        
        prs = Presentation()
        # 16:9 format
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        
        slides_text = response.text.split("SLAYD:")
        for slide_content in slides_text[1:]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = message.text
            slide.placeholders[1].text = slide_content.strip()

        file_path = f"ppt_{message.from_user.id}.pptx"
        prs.save(file_path)
        
        await message.answer_document(FSInputFile(file_path), caption=f"✅ '{message.text}' bo'yicha taqdimot tayyor!")
        os.remove(file_path)
    except Exception as e:
        await message.answer("❌ Taqdimot yaratishda muammo bo'ldi.")
        print(f"ERROR: {e}")
    finally:
        await msg.delete()
        await state.clear()

async def main():
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
